#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Document Structure Extractor (Stage 1)

Produces a clean structure map JSON for:
- PDF
- DOCX
- PPTX
- TXT

No filtering or deduplication happens here.
Stage 1 must remain a faithful representation of the source.
"""

import os
import sys
import json
import argparse
from pathlib import Path

import pdfplumber
import docx
from pptx import Presentation


class DocumentStructureExtractor:

    def __init__(self, input_path: str):
        self.src = Path(input_path)
        if not self.src.exists():
            raise FileNotFoundError(f"File not found: {self.src}")

        self.ext = self.src.suffix.lower()

    # ----------------------------------------------------------
    # Public API
    # ----------------------------------------------------------
    def extract(self) -> dict:
        """Main entry point: determines file type and extracts structure."""
        if self.ext == ".txt":
            return self._extract_txt()
        elif self.ext == ".docx":
            return self._extract_docx()
        elif self.ext == ".pdf":
            return self._extract_pdf()
        elif self.ext == ".pptx":
            return self._extract_pptx()
        else:
            raise ValueError(f"Unsupported file type: {self.ext}")

    def save(self, structure: dict, out_path: Path):
        """Save structure JSON."""
        out_path.parent.mkdir(parents=True, exist_ok=True)
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(structure, f, indent=2, ensure_ascii=False)

    # ----------------------------------------------------------
    # TXT Extractor
    # ----------------------------------------------------------
    def _extract_txt(self) -> dict:
        with open(self.src, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        chunks = []
        for block in text.split("\n\n"):
            block = block.strip()
            if len(block) > 5:
                chunks.append({
                    "page": 1,
                    "chunk_type": "text",
                    "text": block
                })

        return {
            "source_path": str(self.src),
            "doc_type": "txt",
            "total_pages": 1,
            "chunks": chunks,
            "images": []
        }

    # ----------------------------------------------------------
    # DOCX Extractor
    # ----------------------------------------------------------
    def _extract_docx(self) -> dict:
        doc = docx.Document(str(self.src))

        chunks = []
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            if not txt:
                continue

            style = (p.style.name.lower() if p.style else "")
            chunk_type = "heading" if "heading" in style else "text"

            chunks.append({
                "page": 1,
                "chunk_type": chunk_type,
                "text": txt
            })

        return {
            "source_path": str(self.src),
            "doc_type": "docx",
            "total_pages": 1,
            "chunks": chunks,
            "images": []
        }

    # ----------------------------------------------------------
    # PDF Extractor
    # ----------------------------------------------------------
    def _extract_pdf(self) -> dict:
        pdf = pdfplumber.open(str(self.src))

        chunks = []
        images = []

        for page_number, page in enumerate(pdf.pages, start=1):

            # ---- TEXT ----
            raw = page.extract_text() or ""
            raw = raw.replace("\r\n", "\n")
            blocks = [b.strip() for b in raw.split("\n\n") if b.strip()]

            for block in blocks:
                if len(block.split()) > 3:
                    chunks.append({
                        "page": page_number,
                        "chunk_type": "text",
                        "text": block
                    })

            # ---- IMAGES ----
            if page.images:
                for img_index, img in enumerate(page.images):
                    images.append({
                        "page": page_number,
                        "image_index": img_index,
                        "caption": None
                    })

        pdf.close()

        return {
            "source_path": str(self.src),
            "doc_type": "pdf",
            "total_pages": page_number,
            "chunks": chunks,
            "images": images
        }

    # ----------------------------------------------------------
    # PPTX Extractor
    # ----------------------------------------------------------
    def _extract_pptx(self) -> dict:
        prs = Presentation(str(self.src))

        chunks = []
        images = []
        slide_num = 1

        for slide in prs.slides:
            # TEXT
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        chunks.append({
                            "page": slide_num,
                            "chunk_type": "text",
                            "text": txt
                        })

                # IMAGES
                if "picture" in shape.__class__.__name__.lower():
                    images.append({
                        "page": slide_num,
                        "image_index": 0,
                        "caption": None
                    })

            slide_num += 1

        return {
            "source_path": str(self.src),
            "doc_type": "pptx",
            "total_pages": slide_num - 1,
            "chunks": chunks,
            "images": images
        }


# ----------------------------------------------------------
# CLI Entry Point
# ----------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Document Structure Extractor")
    parser.add_argument("input", help="Path to document (.pdf, .docx, .pptx, .txt)")
    parser.add_argument("--out-dir", default="structure_maps", help="Output directory")

    args = parser.parse_args()

    extractor = DocumentStructureExtractor(args.input)
    structure = extractor.extract()

    out_dir = Path(args.out_dir)
    out_path = out_dir / f"{extractor.src.stem}_structure.json"

    extractor.save(structure, out_path)

    print(f"\n[SUCCESS] Structure map written to:\n{out_path}\n")


if __name__ == "__main__":
    main()

# ----------------------------------------------------------
# Main
# ----------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Standalone document structure extractor")
    parser.add_argument("input", help="Path to document (.pdf, .docx, .pptx, .txt)")
    parser.add_argument("--out-dir", default="structure_maps", help="Directory to save structure json")

    args = parser.parse_args()

    src = Path(args.input)
    if not src.exists():
        print(f"[ERROR] File not found: {src}")
        sys.exit(1)

    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    ext = src.suffix.lower()

    if ext == ".txt":
        result = extract_txt(src)
    elif ext == ".docx":
        result = extract_docx(src)
    elif ext == ".pdf":
        result = extract_pdf(src)
    elif ext == ".pptx":
        result = extract_pptx(src)
    else:
        print(f"[ERROR] Unsupported file type: {ext}")
        sys.exit(1)

    out_path = out_dir / f"{src.stem}_structure.json"
    save_json(result, out_path)

    print(f"\n[SUCCESS] Structure map written to:\n{out_path}\n")


if __name__ == "__main__":
    main()
