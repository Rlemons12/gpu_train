#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Document Structure Extractor (Stage 1)

Produces a clean structure map JSON for:
- PDF
- DOCX
- PPTX
- TXT

No filtering or deduplication happens here.
Stage 1 must remain a faithful representation of the source.
"""

import os
import sys
import json
import argparse
from pathlib import Path
import pandas as pd
import pdfplumber
import docx
from pptx import Presentation


class DocumentStructureExtractor:

    def __init__(self, input_path: str):
        self.src = Path(input_path)
        if not self.src.exists():
            raise FileNotFoundError(f"File not found: {self.src}")

        self.ext = self.src.suffix.lower()

        if self.ext == ".doc":
            self.src = self._convert_doc_to_docx(self.src)
            self.ext = ".docx"

    # ----------------------------------------------------------
    # Public API
    # ----------------------------------------------------------
    def extract(self) -> dict:

        if self.ext == ".txt":
            return self._extract_txt()

        elif self.ext == ".docx":
            return self._extract_docx()

        elif self.ext == ".pdf":
            return self._extract_pdf()

        elif self.ext == ".pptx":
            return self._extract_pptx()

        elif self.ext == ".md":
            return self._extract_markdown()

        elif self.ext in {".xlsx", ".xls", ".xlsm"}:
            return self._extract_excel()

        elif self.ext == ".csv":
            return self._extract_csv()

        else:
            raise ValueError(f"Unsupported file type: {self.ext}")

    def save(self, structure: dict, out_path: Path):
        """Save structure JSON."""
        out_path.parent.mkdir(parents=True, exist_ok=True)
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(structure, f, indent=2, ensure_ascii=False)

    # ----------------------------------------------------------
    # TXT Extractor
    # ----------------------------------------------------------
    def _extract_txt(self) -> dict:
        with open(self.src, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        chunks = []
        for block in text.split("\n\n"):
            block = block.strip()
            if len(block) > 5:
                chunks.append({
                    "page": 1,
                    "chunk_type": "text",
                    "text": block
                })

        return {
            "source_path": str(self.src),
            "doc_type": "txt",
            "total_pages": 1,
            "chunks": chunks,
            "images": []
        }

    # ----------------------------------------------------------
    # DOCX Extractor
    # ----------------------------------------------------------
    def _extract_docx(self) -> dict:
        doc = docx.Document(str(self.src))

        chunks = []
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            if not txt:
                continue

            style = (p.style.name.lower() if p.style else "")
            chunk_type = "heading" if "heading" in style else "text"

            chunks.append({
                "page": 1,
                "chunk_type": chunk_type,
                "text": txt
            })

        return {
            "source_path": str(self.src),
            "doc_type": "docx",
            "total_pages": 1,
            "chunks": chunks,
            "images": []
        }

    def _convert_doc_to_docx(self, path: Path) -> Path:
        """
        Convert legacy .doc file to .docx using Microsoft Word COM.
        Windows only.
        """
        try:
            import win32com.client
        except ImportError:
            raise RuntimeError("pywin32 required to convert .doc files")

        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False

        doc = word.Documents.Open(str(path))
        new_path = path.with_suffix(".docx")
        doc.SaveAs(str(new_path), FileFormat=16)  # 16 = wdFormatXMLDocument
        doc.Close()
        word.Quit()

        return new_path


    # ----------------------------------------------------------
    # PDF Extractor
    # ----------------------------------------------------------
    def _extract_pdf(self) -> dict:
        pdf = pdfplumber.open(str(self.src))

        chunks = []
        images = []

        for page_number, page in enumerate(pdf.pages, start=1):

            # ---- TEXT ----
            raw = page.extract_text() or ""
            raw = raw.replace("\r\n", "\n")
            blocks = [b.strip() for b in raw.split("\n\n") if b.strip()]

            for block in blocks:
                if len(block.split()) > 3:
                    chunks.append({
                        "page": page_number,
                        "chunk_type": "text",
                        "text": block
                    })

            # ---- IMAGES ----
            if page.images:
                for img_index, img in enumerate(page.images):
                    images.append({
                        "page": page_number,
                        "image_index": img_index,
                        "caption": None
                    })

        pdf.close()

        return {
            "source_path": str(self.src),
            "doc_type": "pdf",
            "total_pages": page_number,
            "chunks": chunks,
            "images": images
        }

    # ----------------------------------------------------------
    # PPTX Extractor
    # ----------------------------------------------------------
    def _extract_pptx(self) -> dict:
        prs = Presentation(str(self.src))

        chunks = []
        images = []
        slide_num = 1

        for slide in prs.slides:
            # TEXT
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        chunks.append({
                            "page": slide_num,
                            "chunk_type": "text",
                            "text": txt
                        })

                # IMAGES
                if "picture" in shape.__class__.__name__.lower():
                    images.append({
                        "page": slide_num,
                        "image_index": 0,
                        "caption": None
                    })

            slide_num += 1

        return {
            "source_path": str(self.src),
            "doc_type": "pptx",
            "total_pages": slide_num - 1,
            "chunks": chunks,
            "images": images
        }

    def _extract_markdown(self) -> dict:
        chunks = []

        with open(self.src, "r", encoding="utf-8", errors="ignore") as f:
            lines = f.readlines()

        current_section = []
        section_title = "Document Start"
        section_index = 1

        for line in lines:
            stripped = line.strip()

            # New section if heading
            if stripped.startswith("#"):
                if current_section:
                    chunks.append({
                        "page": section_index,
                        "chunk_type": "text",
                        "text": f"{section_title}\n" + "".join(current_section).strip()
                    })
                    section_index += 1
                    current_section = []

                section_title = stripped
            else:
                current_section.append(line)

        # Final section
        if current_section:
            chunks.append({
                "page": section_index,
                "chunk_type": "text",
                "text": f"{section_title}\n" + "".join(current_section).strip()
            })

        return {
            "source_path": str(self.src),
            "doc_type": "md",
            "total_pages": section_index,
            "chunks": chunks,
            "images": []
        }

    def _extract_excel(self) -> dict:
        df = pd.read_excel(self.src, sheet_name=None)

        chunks = []
        page = 1

        for sheet_name, sheet_df in df.items():
            text = sheet_df.astype(str).to_string(index=False)
            chunks.append({
                "page": page,
                "chunk_type": "table",
                "text": f"Sheet: {sheet_name}\n{text}"
            })
            page += 1

        return {
            "source_path": str(self.src),
            "doc_type": self.ext.replace(".", ""),
            "total_pages": page - 1,
            "chunks": chunks,
            "images": []
        }

    def _extract_csv(self):
        import pandas as pd
        df = pd.read_csv(self.src)

        text = df.astype(str).to_string(index=False)

        return {
            "source_path": str(self.src),
            "doc_type": "csv",
            "total_pages": 1,
            "chunks": [{
                "page": 1,
                "chunk_type": "table",
                "text": text
            }],
            "images": []
        }


# ----------------------------------------------------------
# CLI Entry Point
# ----------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description="Document Structure Extractor")
    parser.add_argument("input", help="Path to document (.pdf, .docx, .pptx, .txt)")
    parser.add_argument("--out-dir", default="structure_maps", help="Output directory")

    args = parser.parse_args()

    extractor = DocumentStructureExtractor(args.input)
    structure = extractor.extract()

    out_dir = Path(args.out_dir)
    out_path = out_dir / f"{extractor.src.stem}_structure.json"

    extractor.save(structure, out_path)

    print(f"\n[SUCCESS] Structure map written to:\n{out_path}\n")


if __name__ == "__main__":
    main()


