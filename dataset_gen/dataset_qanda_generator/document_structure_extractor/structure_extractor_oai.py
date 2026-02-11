#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Document Structure Extractor (Stage 1)

Produces a clean structure map JSON for:
- PDF
- DOCX
- PPTX
- TXT
- XLSX / XLSM / XLS

No filtering or deduplication happens here.
Stage 1 must remain a faithful representation of the source.
"""

import sys
import json
import argparse
from pathlib import Path

import pdfplumber
import docx
from pptx import Presentation


class DocumentStructureExtractor:
    MIN_TEXT_WORD_THRESHOLD = 50

    def __init__(self, input_path: str, ocr_engine=None, enable_ocr: bool = True):
        """
        Initialize document extractor.

        Parameters
        ----------
        input_path : str
            Path to document.
        ocr_engine : PaddleOCR | None
            Optional shared OCR engine (recommended for batch processing).
        enable_ocr : bool
            Whether OCR fallback should be available.
        """

        self.src = Path(input_path)

        if not self.src.exists():
            raise FileNotFoundError(f"File not found: {self.src}")

        self.ext = self.src.suffix.lower()

        self.enable_ocr = enable_ocr
        self.ocr_engine = ocr_engine

        # If OCR is enabled but no engine was injected,
        # lazily initialize one.
        if self.enable_ocr and self.ocr_engine is None:
            try:
                from paddleocr import PaddleOCR
                import paddle

                use_gpu = paddle.device.is_compiled_with_cuda()

                self.ocr_engine = PaddleOCR(
                    use_angle_cls=True,
                    use_gpu=use_gpu,
                    lang="en"
                )

                print(f"[OCR INIT] PaddleOCR initialized (GPU={use_gpu})")

            except Exception as e:
                print(f"[OCR INIT FAILED] {e}")
                self.ocr_engine = None
                self.enable_ocr = False

    # ==========================================================
    # PUBLIC API
    # ==========================================================
    def extract(self) -> dict:
        if self.ext == ".txt":
            return self._extract_txt()

        elif self.ext == ".docx":
            return self._extract_docx()

        elif self.ext == ".pdf":
            return self._extract_pdf()

        elif self.ext == ".pptx":
            return self._extract_pptx()

        elif self.ext in [".xlsx", ".xlsm", ".xls"]:
            return self._extract_excel()

        else:
            raise ValueError(f"Unsupported file type: {self.ext}")

    def save(self, structure: dict, out_path: Path):
        out_path.parent.mkdir(parents=True, exist_ok=True)
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(structure, f, indent=2, ensure_ascii=False)

    # ==========================================================
    # TXT
    # ==========================================================
    def _extract_txt(self) -> dict:
        with open(self.src, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        chunks = []
        for block in text.split("\n\n"):
            block = block.strip()
            if len(block) > 5:
                chunks.append({
                    "page": 1,
                    "chunk_type": "text",
                    "text": block
                })

        return {
            "source_path": str(self.src),
            "doc_type": "txt",
            "total_pages": 1,
            "chunks": chunks,
            "images": []
        }

    # ==========================================================
    # DOCX
    # ==========================================================
    def _extract_docx(self) -> dict:
        doc = docx.Document(str(self.src))

        chunks = []
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            if not txt:
                continue

            style = (p.style.name.lower() if p.style else "")
            chunk_type = "heading" if "heading" in style else "text"

            chunks.append({
                "page": 1,
                "chunk_type": chunk_type,
                "text": txt
            })

        return {
            "source_path": str(self.src),
            "doc_type": "docx",
            "total_pages": 1,
            "chunks": chunks,
            "images": []
        }

    # ==========================================================
    # PDF
    # ==========================================================
    def _extract_pdf(self) -> dict:
        from pdf2image import convert_from_path
        import numpy as np

        pdf = pdfplumber.open(str(self.src))

        chunks = []
        images = []
        total_pages = 0

        # -------------------------------
        # NORMAL TEXT EXTRACTION FIRST
        # -------------------------------
        pages_native_text = {}

        for page_number, page in enumerate(pdf.pages, start=1):
            total_pages = page_number

            raw = page.extract_text() or ""
            raw = raw.replace("\r\n", "\n")
            blocks = [b.strip() for b in raw.split("\n\n") if b.strip()]

            page_has_text = False

            for block in blocks:
                if len(block.split()) > 3:
                    page_has_text = True
                    chunks.append({
                        "page": page_number,
                        "chunk_type": "text",
                        "text": block
                    })

            pages_native_text[page_number] = page_has_text

            if page.images:
                for img_index, _ in enumerate(page.images):
                    images.append({
                        "page": page_number,
                        "image_index": img_index,
                        "caption": None
                    })

        pdf.close()

        # --------------------------------
        # OCR FALLBACK (Faithful + Aligned)
        # --------------------------------
        if self.enable_ocr and self.ocr_engine is not None:

            try:
                pages = convert_from_path(str(self.src))

                for page_number, image in enumerate(pages, start=1):

                    # Skip pages that already had usable native text
                    if pages_native_text.get(page_number):
                        continue

                    img_np = np.array(image)
                    result = self.ocr_engine.ocr(img_np, cls=True)

                    page_lines = []

                    if result and result[0]:
                        for line in result[0]:
                            text = line[1][0]
                            confidence = line[1][1]

                            if text.strip() and confidence > 0.5:
                                page_lines.append(text.strip())

                    if page_lines:
                        chunks.append({
                            "page": page_number,
                            "chunk_type": "ocr_text",
                            "text": "\n".join(page_lines)
                        })

                total_pages = len(pages)

            except Exception as e:
                print(f"[OCR ERROR] {self.src} -> {e}")

        return {
            "source_path": str(self.src),
            "doc_type": "pdf",
            "total_pages": total_pages,
            "chunks": chunks,
            "images": images
        }

    # ==========================================================
    # PPTX
    # ==========================================================
    def _extract_pptx(self) -> dict:
        prs = Presentation(str(self.src))

        chunks = []
        images = []
        slide_num = 1

        for slide in prs.slides:

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        chunks.append({
                            "page": slide_num,
                            "chunk_type": "text",
                            "text": txt
                        })

                if "picture" in shape.__class__.__name__.lower():
                    images.append({
                        "page": slide_num,
                        "image_index": 0,
                        "caption": None
                    })

            slide_num += 1

        return {
            "source_path": str(self.src),
            "doc_type": "pptx",
            "total_pages": slide_num - 1,
            "chunks": chunks,
            "images": images
        }

    # ==========================================================
    # EXCEL
    # ==========================================================
    def _extract_excel(self) -> dict:
        from openpyxl import load_workbook

        wb = load_workbook(filename=str(self.src), data_only=True)

        chunks = []
        sheet_count = 0

        for sheet in wb.worksheets:
            sheet_count += 1
            sheet_name = sheet.title

            rows = list(sheet.iter_rows(values_only=True))
            if not rows:
                continue

            headers = [
                str(h).strip() if h else f"Column_{i}"
                for i, h in enumerate(rows[0])
            ]

            for row in rows[1:]:
                row_data = []

                for col_index, cell in enumerate(row):
                    if cell is None:
                        continue

                    header = headers[col_index] if col_index < len(headers) else f"Column_{col_index}"
                    row_data.append(f"{header}: {cell}")

                if not row_data:
                    continue

                text_block = f"Sheet: {sheet_name}\n" + "\n".join(row_data)

                chunks.append({
                    "page": sheet_count,
                    "chunk_type": "table_row",
                    "text": text_block
                })

        return {
            "source_path": str(self.src),
            "doc_type": "excel",
            "total_pages": sheet_count,
            "chunks": chunks,
            "images": []
        }


# ==========================================================
# CLI ENTRY POINT
# ==========================================================
def main():
    parser = argparse.ArgumentParser(description="Document Structure Extractor")
    parser.add_argument("input", help="Path to document")
    parser.add_argument("--out-dir", default="structure_maps", help="Output directory")

    args = parser.parse_args()

    src = Path(args.input)
    if not src.exists():
        print(f"[ERROR] File not found: {src}")
        sys.exit(1)

    try:
        from paddleocr import PaddleOCR
        import paddle

        use_gpu = paddle.device.is_compiled_with_cuda()
        ocr_engine = PaddleOCR(use_angle_cls=True, use_gpu=use_gpu)


    except Exception as e:

        print(f"[OCR INIT FAILED] {e}")

        ocr_engine = None

    extractor = DocumentStructureExtractor(str(src), ocr_engine=ocr_engine)

    structure = extractor.extract()

    out_dir = Path(args.out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"{src.stem}_structure.json"

    extractor.save(structure, out_path)

    print(f"\n[SUCCESS] Structure map written to:\n{out_path}\n")


if __name__ == "__main__":
    main()
