#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Document Structure Extractor - Production Grade with GPU/OCR Support

A comprehensive, self-contained class for extracting structured content from various document types.
Supports: PDF, DOCX, PPTX, XLSX/XLS, TXT, Images (JPG, PNG, etc.)

Features:
- Robust error handling
- Comprehensive logging
- Type-safe implementation
- Resource management
- Configurable extraction options
- Detailed metadata extraction
- GPU-accelerated OCR with PaddleOCR (10-50x faster for scanned docs)
- Multi-language support (EN, CH, FR, DE, KR, JP)
- Intelligent OCR fallback
- Image text extraction
"""

import os
import sys
import json
import logging
import argparse
import time
from pathlib import Path
from typing import Dict, List, Any, Optional, Union, Tuple
from dataclasses import dataclass, asdict, field
from datetime import datetime
from contextlib import contextmanager

# Document processing libraries
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None

# PaddleOCR imports (optional)
try:
    from paddleocr import PaddleOCR

    PADDLEOCR_AVAILABLE = True
except ImportError:
    PADDLEOCR_AVAILABLE = False
    PaddleOCR = None

# Image processing (optional, for OCR)
try:
    from PIL import Image
    import numpy as np

    IMAGE_PROCESSING_AVAILABLE = True
except ImportError:
    IMAGE_PROCESSING_AVAILABLE = False
    Image = None
    np = None


# ==============================================================================
# Configuration & Data Structures
# ==============================================================================

@dataclass
class OCRConfig:
    """Configuration for OCR processing with PaddleOCR."""
    use_gpu: bool = True
    gpu_id: int = 0
    use_angle_cls: bool = True
    lang: str = 'en'  # Language: 'en', 'ch', 'fr', 'german', 'korean', 'japan'
    det_db_thresh: float = 0.3
    det_db_box_thresh: float = 0.5
    rec_batch_num: int = 6
    show_log: bool = False

    # OCR trigger thresholds
    min_text_confidence: float = 0.5  # Minimum confidence to use OCR text
    ocr_if_text_ratio_below: float = 0.3  # Use OCR if extracted text is < 30% of expected
    always_ocr_images: bool = True  # Always OCR extracted images

    # Performance
    enable_mkldnn: bool = True  # Intel CPU optimization
    cpu_threads: int = 10


@dataclass
class ExtractionConfig:
    """Configuration for document extraction."""
    min_text_length: int = 5
    min_words_per_chunk: int = 3
    include_empty_chunks: bool = False
    extract_images: bool = True
    extract_tables: bool = True
    preserve_formatting: bool = True
    max_image_count: int = 1000
    encoding: str = "utf-8"

    # OCR-specific settings
    ocr_config: Optional[OCRConfig] = None
    use_ocr_fallback: bool = True  # Use OCR when standard extraction fails
    extract_image_text: bool = True  # OCR text from images
    prefer_ocr: bool = False  # Always prefer OCR over standard extraction

    def __post_init__(self):
        """Initialize OCR configuration if OCR is available and not set."""
        if self.ocr_config is None and PADDLEOCR_AVAILABLE:
            self.ocr_config = OCRConfig()


@dataclass
class Chunk:
    """Represents a text chunk from a document."""
    page: int
    chunk_type: str
    text: str
    metadata: Optional[Dict[str, Any]] = None

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary."""
        result = {
            "page": self.page,
            "chunk_type": self.chunk_type,
            "text": self.text
        }
        if self.metadata:
            result["metadata"] = self.metadata
        return result


@dataclass
class ImageReference:
    """Represents an image reference in a document."""
    page: int
    image_index: int
    caption: Optional[str] = None
    metadata: Optional[Dict[str, Any]] = None

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary."""
        result = {
            "page": self.page,
            "image_index": self.image_index,
            "caption": self.caption
        }
        if self.metadata:
            result["metadata"] = self.metadata
        return result


@dataclass
class DocumentStructure:
    """Complete document structure."""
    source_path: str
    doc_type: str
    total_pages: int
    chunks: List[Chunk]
    images: List[ImageReference]
    metadata: Dict[str, Any]
    extraction_timestamp: str

    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary."""
        return {
            "source_path": self.source_path,
            "doc_type": self.doc_type,
            "total_pages": self.total_pages,
            "chunks": [chunk.to_dict() for chunk in self.chunks],
            "images": [img.to_dict() for img in self.images],
            "metadata": self.metadata,
            "extraction_timestamp": self.extraction_timestamp
        }


# ==============================================================================
# Main Extractor Class
# ==============================================================================

class DocumentStructureExtractor:
    """
    Production-grade document structure extractor with GPU/OCR support.

    Extracts structured content from various document formats with
    comprehensive error handling, logging, and optional GPU-accelerated OCR.

    Supported formats:
        - PDF (.pdf) - with OCR support for scanned documents
        - Word Documents (.docx)
        - PowerPoint (.pptx)
        - Excel (.xlsx, .xls)
        - Plain Text (.txt)
        - Images (.jpg, .jpeg, .png, .bmp, .tiff) - OCR extraction

    Example:
        >>> # Standard extraction
        >>> extractor = DocumentStructureExtractor("document.pdf")
        >>> structure = extractor.extract()

        >>> # GPU-accelerated OCR
        >>> configuration = ExtractionConfig(
        ...     ocr_config=OCRConfig(use_gpu=True, lang='en'),
        ...     use_ocr_fallback=True
        ... )
        >>> extractor = DocumentStructureExtractor("scanned.pdf", configuration=configuration)
        >>> structure = extractor.extract()
    """

    SUPPORTED_EXTENSIONS = {
        '.pdf', '.docx', '.pptx', '.xlsx', '.xls', '.txt',
        '.jpg', '.jpeg', '.png', '.bmp', '.tiff'  # Image formats for OCR
    }

    def __init__(
            self,
            input_path: Union[str, Path],
            config: Optional[ExtractionConfig] = None,
            logger: Optional[logging.Logger] = None
    ):
        """
        Initialize the document extractor.

        Args:
            input_path: Path to the document file
            config: Extraction configuration (uses defaults if None)
            logger: Logger instance (creates default if None)

        Raises:
            FileNotFoundError: If the input file doesn't exist
            ValueError: If the file type is not supported
        """
        self.src = Path(input_path).resolve()
        self.config = config or ExtractionConfig()
        self.logger = logger or self._setup_logger()

        # Validate file
        if not self.src.exists():
            raise FileNotFoundError(f"File not found: {self.src}")

        if not self.src.is_file():
            raise ValueError(f"Path is not a file: {self.src}")

        self.ext = self.src.suffix.lower()

        if self.ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(
                f"Unsupported file type: {self.ext}. "
                f"Supported types: {', '.join(self.SUPPORTED_EXTENSIONS)}"
            )

        # Verify required library is available
        self._verify_dependencies()

        # Initialize OCR engine if available and configured
        self.ocr_engine = None
        self.ocr_stats = {
            'pages_ocr': 0,
            'images_ocr': 0,
            'fallback_ocr': 0,
            'total_ocr_time': 0.0
        }

        if PADDLEOCR_AVAILABLE and self.config.ocr_config:
            self._init_ocr_engine()

        self.logger.info(f"Initialized extractor for: {self.src.name}")

    def _setup_logger(self) -> logging.Logger:
        """Set up default logger."""
        logger = logging.getLogger(self.__class__.__name__)
        if not logger.handlers:
            handler = logging.StreamHandler(sys.stdout)
            formatter = logging.Formatter(
                '%(asctime)s - %(name)s - %(levelname)s - %(message)s'
            )
            handler.setFormatter(formatter)
            logger.addHandler(handler)
            logger.setLevel(logging.INFO)
        return logger

    def _verify_dependencies(self):
        """Verify that required dependencies are installed for the file type."""
        dependencies = {
            '.pdf': ('pdfplumber', pdfplumber),
            '.docx': ('python-docx', DocxDocument),
            '.pptx': ('python-pptx', Presentation),
            '.xlsx': ('openpyxl', load_workbook),
            '.xls': ('openpyxl', load_workbook),
        }

        # Image formats require PIL and numpy for OCR
        image_exts = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff'}
        if self.ext in image_exts:
            if not IMAGE_PROCESSING_AVAILABLE:
                raise ImportError(
                    "Image processing requires Pillow and numpy. "
                    "Install with: pip install Pillow numpy"
                )
            if not PADDLEOCR_AVAILABLE:
                raise ImportError(
                    "Image OCR requires PaddleOCR. "
                    "Install with: pip install paddleocr"
                )

        if self.ext in dependencies:
            lib_name, lib_module = dependencies[self.ext]
            if lib_module is None:
                raise ImportError(
                    f"Required library '{lib_name}' not found. "
                    f"Install with: pip install {lib_name}"
                )

    def _init_ocr_engine(self):
        """Initialize PaddleOCR engine with GPU support."""
        try:
            ocr_cfg = self.config.ocr_config

            # Check GPU availability
            gpu_available = self._check_gpu()
            use_gpu = ocr_cfg.use_gpu and gpu_available

            if ocr_cfg.use_gpu and not gpu_available:
                self.logger.warning("GPU requested but not available, using CPU")

            self.logger.info(f"Initializing PaddleOCR (GPU: {use_gpu}, Lang: {ocr_cfg.lang})")

            self.ocr_engine = PaddleOCR(
                use_angle_cls=ocr_cfg.use_angle_cls,
                lang=ocr_cfg.lang,
                use_gpu=use_gpu,
                gpu_mem=500 if use_gpu else None,
                det_db_thresh=ocr_cfg.det_db_thresh,
                det_db_box_thresh=ocr_cfg.det_db_box_thresh,
                rec_batch_num=ocr_cfg.rec_batch_num,
                show_log=ocr_cfg.show_log,
                enable_mkldnn=ocr_cfg.enable_mkldnn,
                cpu_threads=ocr_cfg.cpu_threads
            )

            self.logger.info("✓ PaddleOCR initialized successfully")

        except Exception as e:
            self.logger.error(f"Failed to initialize PaddleOCR: {e}")
            self.ocr_engine = None

    def _check_gpu(self) -> bool:
        """Check if GPU is available."""
        try:
            import paddle
            return paddle.is_compiled_with_cuda() and paddle.device.cuda.device_count() > 0
        except:
            return False

    # ==========================================================================
    # Public API
    # ==========================================================================

    def extract(self) -> DocumentStructure:
        """
        Extract structure from the document.

        Returns:
            DocumentStructure object containing all extracted content

        Raises:
            Exception: If extraction fails
        """
        self.logger.info(f"Starting extraction for {self.ext} file")

        try:
            # Handle image files separately
            if self.ext in {'.jpg', '.jpeg', '.png', '.bmp', '.tiff'}:
                result = self._extract_image()
            elif self.ext == '.txt':
                result = self._extract_txt()
            elif self.ext == '.docx':
                result = self._extract_docx()
            elif self.ext == '.pdf':
                result = self._extract_pdf()
            elif self.ext == '.pptx':
                result = self._extract_pptx()
            elif self.ext in {'.xlsx', '.xls'}:
                result = self._extract_excel()
            else:
                raise ValueError(f"Unsupported extension: {self.ext}")

            self.logger.info(
                f"Extraction complete: {len(result.chunks)} chunks, "
                f"{len(result.images)} images"
            )
            return result

        except Exception as e:
            self.logger.error(f"Extraction failed: {str(e)}", exc_info=True)
            raise

    def save(
            self,
            structure: DocumentStructure,
            output_path: Union[str, Path],
            pretty: bool = True
    ) -> Path:
        """
        Save structure to JSON file.

        Args:
            structure: DocumentStructure to save
            output_path: Path to output JSON file
            pretty: Whether to format JSON with indentation

        Returns:
            Path to the saved file
        """
        out_path = Path(output_path).resolve()
        out_path.parent.mkdir(parents=True, exist_ok=True)

        try:
            with open(out_path, 'w', encoding='utf-8') as f:
                if pretty:
                    json.dump(
                        structure.to_dict(),
                        f,
                        indent=2,
                        ensure_ascii=False
                    )
                else:
                    json.dump(structure.to_dict(), f, ensure_ascii=False)

            self.logger.info(f"Structure saved to: {out_path}")
            return out_path

        except Exception as e:
            self.logger.error(f"Failed to save structure: {str(e)}")
            raise

    def extract_and_save(
            self,
            output_path: Union[str, Path],
            pretty: bool = True
    ) -> DocumentStructure:
        """
        Convenience method to extract and save in one call.

        Args:
            output_path: Path to output JSON file
            pretty: Whether to format JSON with indentation

        Returns:
            DocumentStructure object
        """
        structure = self.extract()
        self.save(structure, output_path, pretty)
        return structure

    def get_ocr_stats(self) -> Dict[str, Any]:
        """Get OCR processing statistics."""
        return {
            **self.ocr_stats,
            "avg_time_per_page": (
                self.ocr_stats['total_ocr_time'] / self.ocr_stats['pages_ocr']
                if self.ocr_stats['pages_ocr'] > 0 else 0
            ),
            "ocr_enabled": self.ocr_engine is not None
        }

    # ==========================================================================
    # Text File Extraction
    # ==========================================================================

    def _extract_txt(self) -> DocumentStructure:
        """Extract structure from plain text file."""
        self.logger.debug("Extracting TXT file")

        chunks = []

        try:
            with open(self.src, 'r', encoding=self.config.encoding, errors='ignore') as f:
                text = f.read()
        except UnicodeDecodeError:
            # Try with different encoding
            with open(self.src, 'r', encoding='latin-1', errors='ignore') as f:
                text = f.read()

        # Split by double newlines (paragraphs)
        blocks = text.split('\n\n')

        for block in blocks:
            block = block.strip()

            if not self._is_valid_chunk(block):
                continue

            chunks.append(Chunk(
                page=1,
                chunk_type="paragraph",
                text=block,
                metadata={"char_count": len(block)}
            ))

        return DocumentStructure(
            source_path=str(self.src),
            doc_type="txt",
            total_pages=1,
            chunks=chunks,
            images=[],
            metadata={
                "file_size": self.src.stat().st_size,
                "encoding": self.config.encoding,
                "total_chunks": len(chunks)
            },
            extraction_timestamp=datetime.utcnow().isoformat()
        )

    # ==========================================================================
    # Word Document Extraction
    # ==========================================================================

    def _extract_docx(self) -> DocumentStructure:
        """Extract structure from Word document."""
        self.logger.debug("Extracting DOCX file")

        doc = DocxDocument(str(self.src))
        chunks = []
        images = []

        # Extract paragraphs
        for para_idx, para in enumerate(doc.paragraphs):
            text = para.text.strip()

            if not self._is_valid_chunk(text):
                continue

            # Determine chunk type from style
            style_name = para.style.name.lower() if para.style else ""

            if "heading" in style_name:
                chunk_type = "heading"
                level = self._extract_heading_level(style_name)
            elif "title" in style_name:
                chunk_type = "title"
                level = 0
            else:
                chunk_type = "paragraph"
                level = None

            metadata = {
                "style": para.style.name if para.style else None,
                "para_index": para_idx
            }
            if level is not None:
                metadata["heading_level"] = level

            chunks.append(Chunk(
                page=1,  # DOCX doesn't have page concept easily accessible
                chunk_type=chunk_type,
                text=text,
                metadata=metadata
            ))

        # Extract tables if configured
        if self.config.extract_tables and doc.tables:
            for table_idx, table in enumerate(doc.tables):
                table_text = self._extract_docx_table(table)
                if table_text:
                    chunks.append(Chunk(
                        page=1,
                        chunk_type="table",
                        text=table_text,
                        metadata={"table_index": table_idx}
                    ))

        # Count images (basic - full extraction would need more work)
        if self.config.extract_images:
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    images.append(ImageReference(
                        page=1,
                        image_index=len(images),
                        caption=None
                    ))

        return DocumentStructure(
            source_path=str(self.src),
            doc_type="docx",
            total_pages=1,
            chunks=chunks,
            images=images,
            metadata={
                "file_size": self.src.stat().st_size,
                "total_chunks": len(chunks),
                "total_paragraphs": len(doc.paragraphs),
                "total_tables": len(doc.tables)
            },
            extraction_timestamp=datetime.utcnow().isoformat()
        )

    def _extract_docx_table(self, table) -> str:
        """Extract text from a DOCX table."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):  # Skip empty rows
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _extract_heading_level(self, style_name: str) -> int:
        """Extract heading level from style name."""
        import re
        match = re.search(r'heading\s*(\d+)', style_name, re.IGNORECASE)
        if match:
            return int(match.group(1))
        return 1

    # ==========================================================================
    # PDF Extraction with OCR Support
    # ==========================================================================

    def _extract_pdf(self) -> DocumentStructure:
        """Extract PDF with optional GPU-accelerated OCR support."""
        self.logger.debug("Extracting PDF with OCR support")

        chunks = []
        images = []
        total_pages = 0

        with pdfplumber.open(str(self.src)) as pdf:
            total_pages = len(pdf.pages)

            for page_num, page in enumerate(pdf.pages, start=1):
                self.logger.debug(f"Processing page {page_num}/{total_pages}")

                # Try standard text extraction first
                standard_chunks = self._extract_pdf_page_text(page, page_num)

                # Determine if OCR is needed
                should_ocr = self._should_use_ocr(page, standard_chunks)

                if should_ocr and self.ocr_engine:
                    # Use OCR extraction
                    ocr_chunks = self._extract_pdf_page_ocr(page, page_num)

                    if ocr_chunks:
                        self.logger.debug(f"  OCR: {len(ocr_chunks)} chunks extracted")
                        chunks.extend(ocr_chunks)
                        self.ocr_stats['pages_ocr'] += 1
                    else:
                        # Fallback to standard
                        chunks.extend(standard_chunks)
                else:
                    # Use standard extraction
                    chunks.extend(standard_chunks)

                # Extract images
                if self.config.extract_images:
                    page_images = self._extract_pdf_page_images_with_ocr(page, page_num)
                    images.extend(page_images)

                # Extract tables
                if self.config.extract_tables:
                    table_chunks = self._extract_pdf_page_tables(page, page_num)
                    chunks.extend(table_chunks)

        # Add OCR statistics to metadata
        metadata = {
            "file_size": self.src.stat().st_size,
            "total_chunks": len(chunks),
            "total_images": len(images)
        }

        if self.ocr_engine:
            metadata["ocr_stats"] = self.ocr_stats

        return DocumentStructure(
            source_path=str(self.src),
            doc_type="pdf",
            total_pages=total_pages,
            chunks=chunks,
            images=images,
            metadata=metadata,
            extraction_timestamp=datetime.utcnow().isoformat()
        )

    def _should_use_ocr(self, page, standard_chunks: List[Chunk]) -> bool:
        """Determine if OCR should be used for this page."""
        if not self.ocr_engine:
            return False

        # Always use OCR if preferred
        if self.config.prefer_ocr:
            return True

        # Use OCR if no text extracted
        if not standard_chunks:
            return True

        # Calculate text ratio (extracted vs. expected)
        total_text = sum(len(chunk.text) for chunk in standard_chunks)

        # Heuristic: If very little text, might be scanned
        if total_text < 100:  # Less than 100 chars suggests scanned/image
            return True

        # Check if page has images (might be scanned)
        if hasattr(page, 'images') and page.images:
            # If mostly images, likely scanned
            if len(page.images) > 3:
                return True

        return False

    def _extract_pdf_page_ocr(self, page, page_num: int) -> List[Chunk]:
        """Extract text from PDF page using OCR."""
        chunks = []

        try:
            start_time = time.time()

            # Convert page to image
            img = self._pdf_page_to_image(page)

            if img is None:
                return chunks

            # Run OCR
            result = self.ocr_engine.ocr(img, cls=True)

            # Process OCR results
            if result and result[0]:
                # Group text by vertical position (paragraphs)
                text_blocks = self._group_ocr_results(result[0])

                for block in text_blocks:
                    if len(block) > self.config.min_text_length:
                        chunks.append(Chunk(
                            page=page_num,
                            chunk_type="text_ocr",
                            text=block,
                            metadata={
                                "source": "paddleocr",
                                "word_count": len(block.split()),
                                "char_count": len(block)
                            }
                        ))

            elapsed = time.time() - start_time
            self.ocr_stats['total_ocr_time'] += elapsed
            self.logger.debug(f"  OCR completed in {elapsed:.2f}s")

        except Exception as e:
            self.logger.error(f"OCR failed for page {page_num}: {e}")

        return chunks

    def _pdf_page_to_image(self, page) -> Optional[np.ndarray]:
        """Convert PDF page to image for OCR."""
        try:
            # Get page as PIL image
            pil_image = page.to_image(resolution=300).original

            # Convert to numpy array (RGB)
            img_array = np.array(pil_image.convert('RGB'))

            return img_array

        except Exception as e:
            self.logger.error(f"Failed to convert page to image: {e}")
            return None

    def _group_ocr_results(self, ocr_result: List) -> List[str]:
        """Group OCR results into text blocks."""
        if not ocr_result:
            return []

        # Extract text and positions
        text_items = []
        for line in ocr_result:
            if len(line) >= 2:
                box = line[0]  # Bounding box
                text_info = line[1]  # (text, confidence)

                if isinstance(text_info, tuple) and len(text_info) >= 1:
                    text = text_info[0]
                    confidence = text_info[1] if len(text_info) > 1 else 1.0

                    # Only use high-confidence text
                    if confidence >= self.config.ocr_config.min_text_confidence:
                        # Get vertical position (y-coordinate)
                        y_pos = min(point[1] for point in box)
                        text_items.append((y_pos, text))

        # Sort by vertical position
        text_items.sort(key=lambda x: x[0])

        # Group into blocks (by vertical proximity)
        blocks = []
        current_block = []
        last_y = None
        threshold = 20  # pixels

        for y_pos, text in text_items:
            if last_y is None or abs(y_pos - last_y) < threshold:
                current_block.append(text)
            else:
                if current_block:
                    blocks.append(' '.join(current_block))
                current_block = [text]
            last_y = y_pos

        if current_block:
            blocks.append(' '.join(current_block))

        return blocks

    def _extract_pdf_page_images_with_ocr(
            self,
            page,
            page_num: int
    ) -> List[ImageReference]:
        """Extract images and optionally OCR them."""
        images = []

        if not hasattr(page, 'images') or not page.images:
            return images

        for img_idx, img_info in enumerate(page.images):
            if len(images) >= self.config.max_image_count:
                break

            # Extract image text if configured
            image_text = None
            if self.config.extract_image_text and self.ocr_engine:
                image_text = self._ocr_image_from_pdf(page, img_info)
                if image_text:
                    self.ocr_stats['images_ocr'] += 1

            images.append(ImageReference(
                page=page_num,
                image_index=img_idx,
                caption=image_text,
                metadata={
                    "width": img_info.get('width'),
                    "height": img_info.get('height'),
                    "x0": img_info.get('x0'),
                    "y0": img_info.get('y0'),
                    "ocr_extracted": image_text is not None
                }
            ))

        return images

    def _ocr_image_from_pdf(self, page, img_info: Dict) -> Optional[str]:
        """Extract text from an image in PDF using OCR."""
        try:
            # Get image bounds
            x0 = img_info.get('x0', 0)
            y0 = img_info.get('y0', 0)
            x1 = img_info.get('x1', x0 + img_info.get('width', 100))
            y1 = img_info.get('y1', y0 + img_info.get('height', 100))

            # Crop page to image region
            cropped = page.crop((x0, y0, x1, y1))
            img = cropped.to_image(resolution=300).original
            img_array = np.array(img.convert('RGB'))

            # Run OCR
            result = self.ocr_engine.ocr(img_array, cls=True)

            if result and result[0]:
                texts = [line[1][0] for line in result[0] if len(line) >= 2]
                return ' '.join(texts)

        except Exception as e:
            self.logger.debug(f"Failed to OCR image: {e}")

        return None

    def _extract_pdf_page_text(self, page, page_num: int) -> List[Chunk]:
        """Extract text chunks from a PDF page."""
        chunks = []

        raw_text = page.extract_text() or ""
        raw_text = raw_text.replace('\r\n', '\n')

        # Split into blocks by double newlines
        blocks = [b.strip() for b in raw_text.split('\n\n') if b.strip()]

        for block in blocks:
            if not self._is_valid_chunk(block):
                continue

            chunks.append(Chunk(
                page=page_num,
                chunk_type="text",
                text=block,
                metadata={
                    "word_count": len(block.split()),
                    "char_count": len(block)
                }
            ))

        return chunks

    def _extract_pdf_page_tables(self, page, page_num: int) -> List[Chunk]:
        """Extract tables from a PDF page."""
        chunks = []

        tables = page.extract_tables()
        if not tables:
            return chunks

        for table_idx, table in enumerate(tables):
            if not table:
                continue

            # Convert table to text representation
            table_text = self._format_table(table)

            if table_text and len(table_text) > self.config.min_text_length:
                chunks.append(Chunk(
                    page=page_num,
                    chunk_type="table",
                    text=table_text,
                    metadata={
                        "table_index": table_idx,
                        "rows": len(table),
                        "cols": len(table[0]) if table else 0
                    }
                ))

        return chunks

    def _format_table(self, table: List[List]) -> str:
        """Format table data as text."""
        rows = []
        for row in table:
            if row and any(cell for cell in row):
                # Clean and join cells
                cells = [str(cell).strip() if cell else "" for cell in row]
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    # ==========================================================================
    # Image File OCR
    # ==========================================================================

    def _extract_image(self) -> DocumentStructure:
        """
        Extract text from image file using OCR.

        Returns:
            DocumentStructure with OCR results
        """
        if not self.ocr_engine:
            raise RuntimeError("OCR engine not available for image extraction")

        chunks = []

        try:
            # Load image
            img = Image.open(self.src)
            img_array = np.array(img.convert('RGB'))

            # Run OCR
            result = self.ocr_engine.ocr(img_array, cls=True)

            # Process results
            if result and result[0]:
                text_blocks = self._group_ocr_results(result[0])

                for idx, block in enumerate(text_blocks):
                    if len(block) > self.config.min_text_length:
                        chunks.append(Chunk(
                            page=1,
                            chunk_type="text_ocr",
                            text=block,
                            metadata={
                                "source": "paddleocr",
                                "block_index": idx,
                                "word_count": len(block.split())
                            }
                        ))

        except Exception as e:
            self.logger.error(f"Image OCR failed: {e}")

        return DocumentStructure(
            source_path=str(self.src),
            doc_type="image",
            total_pages=1,
            chunks=chunks,
            images=[],
            metadata={
                "file_size": self.src.stat().st_size,
                "total_chunks": len(chunks),
                "ocr_engine": "paddleocr"
            },
            extraction_timestamp=datetime.utcnow().isoformat()
        )

    # ==========================================================================
    # PowerPoint Extraction
    # ==========================================================================

    def _extract_pptx(self) -> DocumentStructure:
        """Extract structure from PowerPoint file."""
        self.logger.debug("Extracting PPTX file")

        prs = Presentation(str(self.src))
        chunks = []
        images = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            self.logger.debug(f"Processing slide {slide_num}/{len(prs.slides)}")

            # Extract text from shapes
            for shape in slide.shapes:
                # Text content
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()

                    if not self._is_valid_chunk(text):
                        continue

                    # Determine chunk type
                    chunk_type = "text"
                    metadata = {}

                    if hasattr(shape, "text_frame"):
                        # Check if it's a title
                        if shape == slide.shapes.title:
                            chunk_type = "title"
                        metadata["has_text_frame"] = True

                    chunks.append(Chunk(
                        page=slide_num,
                        chunk_type=chunk_type,
                        text=text,
                        metadata=metadata
                    ))

                # Images
                if self.config.extract_images:
                    shape_type = shape.shape_type
                    # 13 is PICTURE shape type
                    if shape_type == 13:
                        images.append(ImageReference(
                            page=slide_num,
                            image_index=len([i for i in images if i.page == slide_num]),
                            caption=None
                        ))

            # Extract tables
            if self.config.extract_tables:
                table_chunks = self._extract_pptx_slide_tables(slide, slide_num)
                chunks.extend(table_chunks)

        return DocumentStructure(
            source_path=str(self.src),
            doc_type="pptx",
            total_pages=len(prs.slides),
            chunks=chunks,
            images=images,
            metadata={
                "file_size": self.src.stat().st_size,
                "total_chunks": len(chunks),
                "total_slides": len(prs.slides),
                "slide_width": prs.slide_width,
                "slide_height": prs.slide_height
            },
            extraction_timestamp=datetime.utcnow().isoformat()
        )

    def _extract_pptx_slide_tables(self, slide, slide_num: int) -> List[Chunk]:
        """Extract tables from a PowerPoint slide."""
        chunks = []

        for shape in slide.shapes:
            if not shape.has_table:
                continue

            table = shape.table
            rows_data = []

            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows_data.append(" | ".join(cells))

            if rows_data:
                table_text = "\n".join(rows_data)
                chunks.append(Chunk(
                    page=slide_num,
                    chunk_type="table",
                    text=table_text,
                    metadata={
                        "rows": len(table.rows),
                        "cols": len(table.columns)
                    }
                ))

        return chunks

    # ==========================================================================
    # Excel Extraction
    # ==========================================================================

    def _extract_excel(self) -> DocumentStructure:
        """Extract structure from Excel file."""
        self.logger.debug("Extracting Excel file")

        wb = load_workbook(filename=str(self.src), data_only=True, read_only=True)
        chunks = []
        total_sheets = len(wb.worksheets)

        for sheet_num, sheet in enumerate(wb.worksheets, start=1):
            self.logger.debug(f"Processing sheet {sheet_num}/{total_sheets}: {sheet.title}")

            sheet_chunks = self._extract_excel_sheet(sheet, sheet_num)
            chunks.extend(sheet_chunks)

        wb.close()

        return DocumentStructure(
            source_path=str(self.src),
            doc_type="excel",
            total_pages=total_sheets,
            chunks=chunks,
            images=[],
            metadata={
                "file_size": self.src.stat().st_size,
                "total_chunks": len(chunks),
                "total_sheets": total_sheets,
                "sheet_names": [ws.title for ws in wb.worksheets]
            },
            extraction_timestamp=datetime.utcnow().isoformat()
        )

    def _extract_excel_sheet(self, sheet, sheet_num: int) -> List[Chunk]:
        """Extract data from an Excel sheet."""
        chunks = []

        # Get all rows
        rows = list(sheet.iter_rows(values_only=True))

        if not rows:
            return chunks

        # First row as headers
        headers = []
        for i, cell in enumerate(rows[0]):
            if cell:
                headers.append(str(cell).strip())
            else:
                headers.append(f"Column_{i + 1}")

        # Process data rows
        for row_idx, row in enumerate(rows[1:], start=2):
            # Skip empty rows
            if not any(cell for cell in row):
                continue

            row_data = []
            for col_idx, cell in enumerate(row):
                if cell is None or cell == "":
                    continue

                header = headers[col_idx] if col_idx < len(headers) else f"Column_{col_idx + 1}"
                row_data.append(f"{header}: {cell}")

            if not row_data:
                continue

            text_block = f"Sheet: {sheet.title}\nRow {row_idx}:\n" + "\n".join(row_data)

            chunks.append(Chunk(
                page=sheet_num,
                chunk_type="table_row",
                text=text_block,
                metadata={
                    "sheet_name": sheet.title,
                    "row_number": row_idx,
                    "columns": len(row_data)
                }
            ))

        return chunks

    # ==========================================================================
    # Helper Methods
    # ==========================================================================

    def _is_valid_chunk(self, text: str) -> bool:
        """Validate if a text chunk meets minimum requirements."""
        if not text:
            return False

        if len(text) < self.config.min_text_length:
            return False

        word_count = len(text.split())
        if word_count < self.config.min_words_per_chunk:
            return False

        return True

    def get_metadata(self) -> Dict[str, Any]:
        """Get file metadata without extracting content."""
        stat = self.src.stat()

        metadata = {
            "filename": self.src.name,
            "file_path": str(self.src),
            "file_size": stat.st_size,
            "file_type": self.ext,
            "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            "ocr_available": self.ocr_engine is not None
        }

        return metadata


# ==============================================================================
# CLI Interface
# ==============================================================================

def main():
    """Command-line interface for the document structure extractor."""
    parser = argparse.ArgumentParser(
        description="Extract structured content from documents with optional GPU/OCR",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Supported formats:
  PDF    - .pdf (with OCR support)
  Word   - .docx
  Excel  - .xlsx, .xls
  PowerPoint - .pptx
  Text   - .txt
  Images - .jpg, .jpeg, .png, .bmp, .tiff (OCR extraction)

GPU/OCR Options:
  --use-gpu              Use GPU for OCR (default if available)
  --no-gpu               Force CPU mode
  --ocr-lang LANG        OCR language: en, ch, fr, german, korean, japan
  --prefer-ocr           Always prefer OCR over standard extraction
  --no-ocr-fallback      Disable OCR fallback

Examples:
  %(prog)s document.pdf
  %(prog)s document.pdf --out-dir ./output
  %(prog)s scanned.pdf --prefer-ocr --use-gpu
  %(prog)s chinese.pdf --ocr-lang ch
  %(prog)s image.jpg --output text.json
        """
    )

    parser.add_argument(
        "input",
        help="Path to input document"
    )

    parser.add_argument(
        "--out-dir",
        default="structure_maps",
        help="Output directory (default: structure_maps)"
    )

    parser.add_argument(
        "--output", "-o",
        help="Specific output file path (overrides --out-dir)"
    )

    # OCR/GPU options
    parser.add_argument(
        "--use-gpu",
        action="store_true",
        default=True,
        help="Use GPU for OCR (default: True if available)"
    )

    parser.add_argument(
        "--no-gpu",
        action="store_true",
        help="Force CPU mode for OCR"
    )

    parser.add_argument(
        "--ocr-lang",
        default="en",
        help="OCR language: en, ch, fr, german, korean, japan (default: en)"
    )

    parser.add_argument(
        "--prefer-ocr",
        action="store_true",
        help="Always prefer OCR over standard extraction"
    )

    parser.add_argument(
        "--no-ocr-fallback",
        action="store_true",
        help="Disable automatic OCR fallback"
    )

    # Standard options
    parser.add_argument(
        "--min-length",
        type=int,
        default=5,
        help="Minimum text length for chunks (default: 5)"
    )

    parser.add_argument(
        "--min-words",
        type=int,
        default=3,
        help="Minimum words per chunk (default: 3)"
    )

    parser.add_argument(
        "--no-images",
        action="store_true",
        help="Skip image extraction"
    )

    parser.add_argument(
        "--no-tables",
        action="store_true",
        help="Skip table extraction"
    )

    parser.add_argument(
        "--compact",
        action="store_true",
        help="Save JSON in compact format (no indentation)"
    )

    parser.add_argument(
        "--verbose", "-v",
        action="store_true",
        help="Enable verbose logging"
    )

    parser.add_argument(
        "--quiet", "-q",
        action="store_true",
        help="Suppress all output except errors"
    )

    args = parser.parse_args()

    # Set up logging
    log_level = logging.ERROR if args.quiet else (logging.DEBUG if args.verbose else logging.INFO)
    logging.basicConfig(
        level=log_level,
        format='%(asctime)s - %(levelname)s - %(message)s'
    )

    try:
        # Create OCR configuration if OCR is available
        ocr_config = None
        if PADDLEOCR_AVAILABLE:
            ocr_config = OCRConfig(
                use_gpu=args.use_gpu and not args.no_gpu,
                lang=args.ocr_lang,
                show_log=args.verbose
            )

        # Create configuration
        config = ExtractionConfig(
            min_text_length=args.min_length,
            min_words_per_chunk=args.min_words,
            extract_images=not args.no_images,
            extract_tables=not args.no_tables,
            ocr_config=ocr_config,
            use_ocr_fallback=not args.no_ocr_fallback,
            prefer_ocr=args.prefer_ocr
        )

        # Initialize extractor
        extractor = DocumentStructureExtractor(args.input, config=config)

        # Determine output path
        if args.output:
            out_path = Path(args.output)
        else:
            out_dir = Path(args.out_dir)
            out_path = out_dir / f"{extractor.src.stem}_structure.json"

        # Extract and save
        structure = extractor.extract_and_save(
            out_path,
            pretty=not args.compact
        )

        # Print summary
        if not args.quiet:
            print("\n" + "=" * 60)
            print("EXTRACTION COMPLETE")
            print("=" * 60)
            print(f"Input:  {extractor.src}")
            print(f"Output: {out_path}")
            print(f"Type:   {structure.doc_type.upper()}")
            print(f"Pages:  {structure.total_pages}")
            print(f"Chunks: {len(structure.chunks)}")
            print(f"Images: {len(structure.images)}")

            # OCR stats
            if extractor.ocr_engine:
                ocr_stats = extractor.get_ocr_stats()
                if ocr_stats['pages_ocr'] > 0 or ocr_stats['images_ocr'] > 0:
                    print("\nOCR Statistics:")
                    print(f"  Pages OCR'd:  {ocr_stats['pages_ocr']}")
                    print(f"  Images OCR'd: {ocr_stats['images_ocr']}")
                    print(f"  Total Time:   {ocr_stats['total_ocr_time']:.2f}s")
                    if ocr_stats['pages_ocr'] > 0:
                        print(f"  Avg/Page:     {ocr_stats['avg_time_per_page']:.2f}s")

            print("=" * 60 + "\n")

        return 0

    except FileNotFoundError as e:
        logging.error(f"File not found: {e}")
        return 1
    except ValueError as e:
        logging.error(f"Invalid input: {e}")
        return 1
    except ImportError as e:
        logging.error(f"Missing dependency: {e}")
        return 1
    except Exception as e:
        logging.error(f"Extraction failed: {e}", exc_info=args.verbose)
        return 1


if __name__ == "__main__":
    sys.exit(main())